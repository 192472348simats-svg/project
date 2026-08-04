import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE


def build_default_template(output_path: str = None) -> str:
    """
    Creates a default mentor_template.pptx that matches the exact visual structure,
    header banners, tables, and parent letter layouts of NEW.pptx.
    """
    if not output_path:
        base_dir = os.path.dirname(os.path.abspath(__file__))
        output_path = os.path.join(base_dir, "mentor_template.pptx")

    os.makedirs(os.path.dirname(output_path), exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]

    # -------------------------------------------------------------
    # SLIDE 1: STUDENT PROFILE & COURSE SUMMARY TABLE
    # -------------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)

    # Top Banner Header
    banner1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.2))
    banner1.fill.solid()
    banner1.fill.fore_color.rgb = RGBColor(120, 0, 0) # Dark burgundy red background
    banner1.line.color.rgb = RGBColor(120, 0, 0)

    tf1 = banner1.text_frame
    tf1.margin_left = Inches(0.8)
    tf1.margin_top = Inches(0.2)
    p1 = tf1.paragraphs[0]
    p1.text = "SIMATS ENGINEERING"
    p1.font.size = Pt(36)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(255, 255, 255)

    # Left Box: Student Name & Photo Box
    name_box = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.5), Inches(3.8), Inches(0.5))
    name_box.fill.solid()
    name_box.fill.fore_color.rgb = RGBColor(241, 245, 249)
    name_box.line.color.rgb = RGBColor(203, 213, 225)

    tf_name = name_box.text_frame
    p_name = tf_name.paragraphs[0]
    p_name.text = "ARUNPRASATH R, 192524077"
    p_name.font.size = Pt(13)
    p_name.font.bold = True
    p_name.alignment = PP_ALIGN.CENTER
    p_name.font.color.rgb = RGBColor(15, 23, 42)

    # Photo Box placeholder
    photo_box = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(2.1), Inches(3.8), Inches(4.8))
    photo_box.fill.solid()
    photo_box.fill.fore_color.rgb = RGBColor(226, 232, 240)
    photo_box.line.color.rgb = RGBColor(203, 213, 225)

    # Right Box: Academic Details Table
    rows = 9
    cols = 4
    left = Inches(4.7)
    top = Inches(1.5)
    width = Inches(8.0)
    height = Inches(5.4)

    table_shape = slide1.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    table.columns[0].width = Inches(1.0) # SLOT
    table.columns[1].width = Inches(1.6) # SUB CODE
    table.columns[2].width = Inches(4.2) # SUB NAME
    table.columns[3].width = Inches(1.2) # ATTD %

    # Header Row: OVERALL DETAILS
    cell_hdr = table.cell(0, 0)
    cell_hdr.merge(table.cell(0, 3))
    cell_hdr.text = "OVERALL DETAILS"
    cell_hdr.fill.solid()
    cell_hdr.fill.fore_color.rgb = RGBColor(234, 113, 35) # Orange header
    p = cell_hdr.text_frame.paragraphs[0]
    p.font.bold = True
    p.font.size = Pt(12)
    p.alignment = PP_ALIGN.CENTER
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Row 1: NUMBER OF COURSES COMPLETED
    c1 = table.cell(1, 0)
    c1.merge(table.cell(1, 3))
    c1.text = "NUMBER OF COURSES COMPLETED : 16"
    p = c1.text_frame.paragraphs[0]
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Row 2: NUMBER OF ARREARS
    c2 = table.cell(2, 0)
    c2.merge(table.cell(2, 3))
    c2.text = "NUMBER OF ARREARS : 01"
    p = c2.text_frame.paragraphs[0]
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Row 3: Blank separator
    c3 = table.cell(3, 0)
    c3.merge(table.cell(3, 3))
    c3.text = ""

    # Row 4: DETAILS ABOUT CURRENTLY DOING COURSE
    c4 = table.cell(4, 0)
    c4.merge(table.cell(4, 3))
    c4.text = "DETAILS ABOUT CURRENTLY DOING COURSE"
    p = c4.text_frame.paragraphs[0]
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Row 5: Table Column Titles
    table.cell(5, 0).text = "SLOT"
    table.cell(5, 1).text = "SUB CODE"
    table.cell(5, 2).text = "SUB NAME"
    table.cell(5, 3).text = "ATTD %"

    for c in range(4):
        p = table.cell(5, c).text_frame.paragraphs[0]
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

    # Rows 6-8: Slots A, B, C, D
    slots_data = [
        ("A", "CSA0510", "Database Management System", "82%"),
        ("B", "CSA1415", "Compiler Design", "80%"),
        ("C", "CSA1603", "Data Warehousing and Data Mining", "80%"),
        ("D", "CSA0314", "Data Structures", "78%")
    ]

    for idx, (slot, code, sname, att) in enumerate(slots_data[:3]):
        r = 6 + idx
        table.cell(r, 0).text = slot
        table.cell(r, 1).text = code
        table.cell(r, 2).text = sname
        table.cell(r, 3).text = att
        table.cell(r, 0).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        table.cell(r, 3).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Footer Date
    tx_foot1 = slide1.shapes.add_textbox(Inches(10.5), Inches(7.0), Inches(2.5), Inches(0.4))
    tx_foot1.text_frame.text = "GP July 10th to 16th 2026"
    tx_foot1.text_frame.paragraphs[0].font.size = Pt(9)
    tx_foot1.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    # -------------------------------------------------------------
    # SLIDE 2: PARENT LETTER / MODULAR REPORT
    # -------------------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)

    # Banner
    banner2 = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.2))
    banner2.fill.solid()
    banner2.fill.fore_color.rgb = RGBColor(120, 0, 0)
    banner2.line.color.rgb = RGBColor(120, 0, 0)

    tf2 = banner2.text_frame
    tf2.margin_left = Inches(0.8)
    tf2.margin_top = Inches(0.2)
    p2 = tf2.paragraphs[0]
    p2.text = "SIMATS ENGINEERING"
    p2.font.size = Pt(36)
    p2.font.bold = True
    p2.font.color.rgb = RGBColor(255, 255, 255)

    # Main Letter Text Box
    tx_letter = slide2.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(12.1), Inches(5.6))
    tf_let = tx_letter.text_frame
    tf_let.word_wrap = True

    p_let = tf_let.paragraphs[0]
    p_let.text = (
        "Dear Parents, The previous slide contains a detailed summary of your ward's academic progress, "
        "including the subjects he has successfully completed, the number of pending arrears (if any), "
        "the courses he is currently pursuing, and his attendance percentage in each subject. We kindly request "
        "you to spend a few minutes reviewing these details carefully. A constructive discussion with your son "
        "regarding his academic strengths, areas that require improvement, and future learning goals will greatly motivate him to perform better.\n\n"
        "We are pleased to inform you that our institution successfully celebrated SIMMAM 2026, the College Cultural Festival, "
        "providing students with an excellent platform to exhibit their creativity, talents, leadership qualities, and team spirit.\n\n"
        "Regarding his NPTEL enrolment, he has registered for the following courses:\n"
        "•  Ethical Hacking – The course is scheduled to commence on 20 July.\n"
        "•  Management Information System – The course is scheduled to commence on 20 July.\n\n"
        "We also request your support in reminding him to strictly follow the institutional code of conduct. He should wear his official "
        "SIMATS Identity Card at all times while inside the campus. His hair should be neatly maintained, and he should either remain clean-shaven "
        "or keep his beard properly trimmed. In addition to regular academics, we strongly encourage him to participate actively in technical events "
        "such as hackathons, seminars, workshops, coding competitions, and code debugging activities. With Regards, Dr. T. Kumaragurubaran, 7373032383."
    )
    p_let.font.size = Pt(11)
    p_let.font.color.rgb = RGBColor(255, 255, 255)

    # Footer Date
    tx_foot2 = slide2.shapes.add_textbox(Inches(10.5), Inches(7.0), Inches(2.5), Inches(0.4))
    tx_foot2.text_frame.text = "GP July 10th to 16th 2026"
    tx_foot2.text_frame.paragraphs[0].font.size = Pt(9)
    tx_foot2.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    prs.save(output_path)
    return output_path


if __name__ == "__main__":
    path = build_default_template()
    print(f"NEW.pptx style default mentor template built successfully at: {path}")
