import copy
import re
from typing import Dict, Any, List, Optional
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches, Pt
from .analytics import StudentAnalytics
from .paragraph_builder import ParagraphBuilder


class TemplateMapper:
    """
    Analyzes uploaded PowerPoint templates, processes double-brace placeholders
    (e.g. {{NAME}}, {{ROLL}}, {{PHOTO}}, {{ATTENDANCE_TABLE}}, {{PARENT_REPORT}},
    {{RECENT_EVENTS}}, {{NPTEL}}, {{MENTOR_NAME}}, {{MENTOR_PHONE}}), and populates
    shapes, tables, and images dynamically while preserving 100% of original visual design.
    """

    @classmethod
    def scan_template_structure(cls, prs: Presentation) -> Dict[str, Any]:
        """
        Inspects slide 1 and slide 2 of template to locate dynamic shapes.
        """
        if len(prs.slides) < 2:
            return {"valid": False, "reason": "Template must contain at least 2 slides (Slide 1: Profile/Table, Slide 2: Parent Letter)."}

        slide1 = prs.slides[0]
        slide2 = prs.slides[1]

        s1_table = None
        s1_name_shape = None
        s1_photo_shape = None

        for shape in slide1.shapes:
            if shape.has_table:
                s1_table = shape.table
            elif shape.has_text_frame:
                txt = shape.text_frame.text.upper()
                if any(k in txt for k in ["{{NAME}}", "{{ROLL}}", "1925", "REG", ",", "ARUNPRASATH", "NAME"]):
                    s1_name_shape = shape
            elif shape.shape_type in [MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.RECTANGLE]:
                if shape.width > Inches(1.5) and shape.height > Inches(1.5):
                    s1_photo_shape = shape

        s2_letter_shape = None
        for shape in slide2.shapes:
            if shape.has_text_frame:
                txt = shape.text_frame.text
                if any(k in txt for k in ["{{PARENT_REPORT}}", "{{PARENT_LETTER}}", "Dear Parents", "SIMMAM", "NPTEL"]):
                    s2_letter_shape = shape
                    break

        return {
            "valid": True,
            "s1_table": s1_table,
            "s1_name_shape": s1_name_shape,
            "s1_photo_shape": s1_photo_shape,
            "s2_letter_shape": s2_letter_shape
        }

    @classmethod
    def process_slide_placeholders(
        cls,
        slide,
        analytics: StudentAnalytics,
        student_record: Dict[str, Any],
        recent_events: str = "SIMMAM 2026",
        mentor_name: str = "Dr. T. Kumaragurubaran",
        mentor_phone: str = "7373032383",
        photo_path: Optional[str] = None
    ) -> bool:
        """
        Scans all shapes on a slide and replaces explicit {{PLACEHOLDER}} tags dynamically.
        Returns True if any explicit {{...}} tag was found and processed.
        """
        parent_letter_text = ParagraphBuilder.generate_parent_letter(
            analytics,
            recent_events=recent_events,
            mentor_name=mentor_name,
            mentor_phone=mentor_phone
        )

        replacements = {
            "{{NAME}}": analytics.name.upper(),
            "{{STUDENT_NAME}}": analytics.name.upper(),
            "{{ROLL}}": analytics.reg_no,
            "{{REG_NO}}": analytics.reg_no,
            "{{REGNO}}": analytics.reg_no,
            "{{DEPARTMENT}}": analytics.department,
            "{{DEPT}}": analytics.department,
            "{{OVERALL_ATTENDANCE}}": f"{analytics.overall_attendance:.1f}%",
            "{{ARREARS}}": "NIL" if analytics.arrears == 0 else f"{analytics.arrears:02d}",
            "{{COMPLETED_COURSES}}": str(analytics.completed_courses),
            "{{RECENT_EVENTS}}": recent_events,
            "{{MENTOR_NAME}}": mentor_name,
            "{{MENTOR_PHONE}}": mentor_phone,
            "{{NPTEL}}": analytics.nptel_courses,
            "{{PARENT_REPORT}}": parent_letter_text,
            "{{PARENT_LETTER}}": parent_letter_text,
        }

        tag_found = False
        photo_placed = False

        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                for r in table.rows:
                    for c in r.cells:
                        if "{{ATTENDANCE_TABLE}}" in c.text:
                            tag_found = True
                            c.text = c.text.replace("{{ATTENDANCE_TABLE}}", "")
                cls.populate_slide_1_table(table, analytics, student_record)

            elif shape.has_text_frame:
                txt = shape.text_frame.text
                if "{{" in txt and "}}" in txt:
                    tag_found = True

                if "{{PHOTO}}" in txt:
                    if photo_path:
                        photo_placed = cls._replace_shape_with_photo(slide, shape, photo_path) or photo_placed
                    continue

                if "{{PARENT_REPORT}}" in txt or "{{PARENT_LETTER}}" in txt:
                    shape.text_frame.word_wrap = True
                    cls._replace_text_in_frame(shape.text_frame, "{{PARENT_REPORT}}", parent_letter_text)
                    cls._replace_text_in_frame(shape.text_frame, "{{PARENT_LETTER}}", parent_letter_text)
                    continue

                for tag, val in replacements.items():
                    if tag in txt:
                        cls._replace_text_in_frame(shape.text_frame, tag, val)

        return {"tag_found": tag_found, "photo_placed": photo_placed}

    @classmethod
    def _replace_text_in_frame(cls, text_frame, search_text: str, replace_text: str):
        if search_text not in text_frame.text:
            return

        for p in text_frame.paragraphs:
            if search_text in p.text:
                if p.runs:
                    full_text = "".join(r.text for r in p.runs)
                    if search_text in full_text:
                        new_text = full_text.replace(search_text, replace_text)
                        p.runs[0].text = new_text
                        for r in p.runs[1:]:
                            r.text = ""
                    else:
                        p.text = p.text.replace(search_text, replace_text)
                else:
                    p.text = p.text.replace(search_text, replace_text)

    @classmethod
    def _replace_shape_with_photo(cls, slide, shape, photo_path: str) -> bool:
        """
        Inserts the student's photo at the exact position/size of the given
        placeholder shape (preserving the mentor's original layout), then
        clears the placeholder's own text so no tag text leaks through.
        Returns True if the photo was successfully inserted.
        """
        try:
            left, top, width, height = shape.left, shape.top, shape.width, shape.height
            if shape.has_text_frame:
                shape.text_frame.text = ""
            slide.shapes.add_picture(photo_path, left, top, width=width, height=height)
            return True
        except Exception:
            return False

    @classmethod
    def find_generic_photo_placeholder(cls, slide):
        """
        Heuristically locates a mentor-provided photo placeholder shape when no
        explicit {{PHOTO}} tag is present in the template. Looks for an empty
        (no meaningful text), reasonably square-to-portrait shape of a size
        typical for a photo slot, and returns the largest such candidate.
        Returns None if no suitable shape is found.
        """
        candidates = []
        for shape in slide.shapes:
            try:
                width, height = shape.width, shape.height
            except Exception:
                continue
            if width is None or height is None:
                continue
            if width < Inches(1.2) or height < Inches(1.2):
                continue
            if shape.has_table:
                continue
            has_text = shape.has_text_frame and shape.text_frame.text.strip() != ""
            if has_text:
                continue

            aspect = width / height
            if 0.5 <= aspect <= 1.6:
                candidates.append(shape)

        if not candidates:
            return None

        candidates.sort(key=lambda s: (s.width * s.height), reverse=True)
        return candidates[0]

    @classmethod
    def populate_slide_1_table(cls, table, analytics: StudentAnalytics, student_record: Dict[str, Any]):
        """
        Populates Slide 1 table matching NEW.pptx layout safely and dynamically:
        - OVERALL DETAILS: Completed Courses & Arrears
        - DETAILS ABOUT CURRENTLY DOING COURSE: Slots A, B, C, D (SUB CODE, SUB NAME, ATTD %)
        """
        if not table or not hasattr(table, "rows") or len(table.rows) == 0:
            return

        subject_att = student_record.get("subject_attendance", {})
        subj_list = list(subject_att.items()) if subject_att else []
        slots = ["A", "B", "C", "D"]
        slot_idx = 0

        for r_idx, row in enumerate(table.rows):
            num_cells = len(row.cells)
            if num_cells == 0:
                continue

            cell_texts = [row.cells[c].text.strip().upper() for c in range(num_cells)]
            row_str = " ".join(cell_texts)

            if "NUMBER OF COURSES COMPLETED" in row_str:
                cls._update_summary_row(row, "NUMBER OF COURSES COMPLETED", str(analytics.completed_courses))

            elif "NUMBER OF ARREARS" in row_str:
                arr_str = "NIL" if analytics.arrears == 0 else f"{analytics.arrears:02d}"
                cls._update_summary_row(row, "NUMBER OF ARREARS", arr_str)

            elif num_cells >= 2 and cell_texts[0] in slots:
                if slot_idx < len(subj_list):
                    code_name, att_val = subj_list[slot_idx]
                    if " " in code_name:
                        parts = code_name.split(" ", 1)
                        sub_code = parts[0]
                        sub_name = parts[1]
                    else:
                        sub_code = code_name
                        sub_name = code_name

                    att_str = f"{att_val:.0f}%" if att_val > 0 else f"{att_val:.0f}"

                    if num_cells >= 4:
                        cls._set_table_cell_text(row.cells[1], sub_code)
                        cls._set_table_cell_text(row.cells[2], sub_name)
                        cls._set_table_cell_text(row.cells[3], att_str)
                    elif num_cells == 3:
                        cls._set_table_cell_text(row.cells[1], f"{sub_code} - {sub_name}")
                        cls._set_table_cell_text(row.cells[2], att_str)
                    elif num_cells == 2:
                        cls._set_table_cell_text(row.cells[1], f"{sub_code} {sub_name} ({att_str})")
                else:
                    if num_cells >= 4:
                        cls._set_table_cell_text(row.cells[1], "")
                        cls._set_table_cell_text(row.cells[2], "YET TO ALLOCATE")
                        cls._set_table_cell_text(row.cells[3], "")
                    elif num_cells >= 2:
                        cls._set_table_cell_text(row.cells[1], "YET TO ALLOCATE")

                slot_idx += 1

    @classmethod
    def _update_summary_row(cls, row, keyword: str, new_value: str):
        num_cells = len(row.cells)
        if num_cells == 0:
            return

        for c_idx in range(num_cells):
            cell_txt = row.cells[c_idx].text.strip()
            if keyword in cell_txt.upper():
                if ":" in cell_txt:
                    prefix = cell_txt.split(":")[0].strip()
                    cls._set_table_cell_text(row.cells[c_idx], f"{prefix} : {new_value}")
                    return
                elif num_cells > 1:
                    last_cell_idx = num_cells - 1
                    cls._set_table_cell_text(row.cells[last_cell_idx], new_value)
                    return
                else:
                    cls._set_table_cell_text(row.cells[c_idx], f"{keyword} : {new_value}")
                    return

        if num_cells > 0:
            cls._set_table_cell_text(row.cells[num_cells - 1], new_value)

    @staticmethod
    def _set_table_cell_text(cell, new_text: str):
        if not hasattr(cell, "text_frame") or cell.text_frame is None:
            return

        if not cell.text_frame.paragraphs:
            cell.text = new_text
            return

        p = cell.text_frame.paragraphs[0]
        if p.runs:
            p.runs[0].text = new_text
            for r in p.runs[1:]:
                r.text = ""
        else:
            p.text = new_text