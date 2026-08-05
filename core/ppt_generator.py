import os
import sys
import io
import copy
import shutil
import tempfile
import traceback
from typing import List, Optional, Callable
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE

from .analytics import StudentAnalytics
from .recommendations import RecommendationEngine
from .paragraph_builder import ParagraphBuilder
from .template_mapper import TemplateMapper
from .image_handler import ImageHandler


class PPTGenerationError(Exception):
    """
    Detailed exception class for PowerPoint generation errors.
    Captures file, function, student, slide, table, reason, and full traceback.
    """
    def __init__(
        self,
        file_name: str,
        function_name: str,
        student_name: str,
        slide_name: str,
        table_name: str,
        reason: str,
        traceback_str: str = ""
    ):
        self.file_name = file_name
        self.function_name = function_name
        self.student_name = student_name
        self.slide_name = slide_name
        self.table_name = table_name
        self.reason = reason
        self.traceback_str = traceback_str
        super().__init__(self.formatted_message())

    def formatted_message(self) -> str:
        msg = (
            f"PowerPoint Generation Failed\n\n"
            f"File:\n{self.file_name}\n\n"
            f"Function:\n{self.function_name}\n\n"
            f"Student:\n{self.student_name}\n\n"
            f"Slide:\n{self.slide_name}\n\n"
            f"Table:\n{self.table_name}\n\n"
            f"Reason:\n{self.reason}"
        )
        if self.traceback_str:
            msg += f"\n\nTraceback:\n{self.traceback_str}"
        return msg


class PPTReportGenerator:
    """
    Generates PowerPoint presentations (.pptx) matching the exact visual layout of NEW.pptx.
    Duplicates template slides while preserving original background, fonts, borders, headers, and colors.
    """

    def __init__(
        self,
        template_path: Optional[str] = None,
        image_folder: Optional[str] = None,
        recent_events: str = "SIMMAM 2026",
        mentor_name: str = "Dr. T. Kumaragurubaran",
        mentor_phone: str = "7373032383"
    ):
        self.template_path = template_path if (template_path and os.path.exists(template_path)) else None
        self.image_handler = ImageHandler(image_folder)
        self.recent_events = recent_events
        self.mentor_name = mentor_name
        self.mentor_phone = mentor_phone

    def generate_all_reports(
        self,
        analytics_list: List[StudentAnalytics],
        export_dir: str,
        student_records: Optional[List[dict]] = None,
        progress_callback: Optional[Callable[[int, int, str], None]] = None
    ) -> dict:
        """
        Generates individual PPTX files and one combined presentation PPTX file.
        """
        os.makedirs(export_dir, exist_ok=True)
        indiv_dir = os.path.join(export_dir, "individual_reports")
        os.makedirs(indiv_dir, exist_ok=True)

        combined_pptx_path = os.path.join(export_dir, "Combined_Cohort_Mentor_Report.pptx")
        individual_files = []

        total = len(analytics_list)
        if total == 0:
            raise ValueError("No student records provided for PowerPoint generation.")

        if progress_callback:
            progress_callback(0, total, "Reading Template...")

        for i, analytics in enumerate(analytics_list):
            if progress_callback:
                progress_callback(i + 1, total, f"Generating report for {analytics.name} ({analytics.reg_no})...")

            student_rec = student_records[i] if (student_records and i < len(student_records)) else {}

            try:
                indiv_prs = self._create_student_presentation(analytics, student_rec)
                indiv_path = os.path.join(indiv_dir, f"Mentor_Report_{analytics.reg_no}.pptx")
                indiv_prs.save(indiv_path)
                individual_files.append(indiv_path)
            except PPTGenerationError:
                raise
            except Exception as e:
                tb_str = traceback.format_exc()
                print("\n" + "="*80, file=sys.stderr)
                print("PPT GENERATION EXCEPTION DETECTED:", file=sys.stderr)
                print(tb_str, file=sys.stderr)
                print("="*80 + "\n", file=sys.stderr)

                tb_lines = [line.strip() for line in tb_str.splitlines() if line.strip()]
                file_info = "ppt_generator.py"
                func_info = "generate_all_reports()"
                for line in reversed(tb_lines):
                    if "File " in line and "in " in line:
                        parts = line.split(",")
                        file_info = parts[0].replace('File "', '').replace('"', '').strip()
                        if len(parts) >= 3:
                            func_info = parts[2].strip()
                        break

                raise PPTGenerationError(
                    file_name=file_info,
                    function_name=func_info,
                    student_name=analytics.name,
                    slide_name="Slide 1 (Profile & Summary)",
                    table_name="Attendance & Academic Table",
                    reason=str(e),
                    traceback_str=tb_str
                ) from e

        if progress_callback:
            progress_callback(total, total, "Combining Cohort Presentations...")

        self._build_combined_presentation(individual_files, combined_pptx_path)

        if progress_callback:
            progress_callback(total, total, "Export Complete")

        return {
            "combined_ppt": combined_pptx_path,
            "individual_ppts": individual_files,
            "total_generated": len(individual_files)
        }

    def _load_template_prs(self) -> Presentation:
        if self.template_path and os.path.exists(self.template_path):
            return Presentation(self.template_path)
        
        default_path = os.path.abspath("templates/mentor_template.pptx")
        if os.path.exists(default_path):
            return Presentation(default_path)
        
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        return prs

    @staticmethod
    def _trim_to_template_slides(prs: Presentation, num_template_slides: int = 2):
        """
        Trims loaded presentation to the base template slide count (2 slides),
        preventing slide count explosion when loading previously generated template files.
        """
        sldIdLst = prs.slides._sldIdLst
        while len(sldIdLst) > num_template_slides:
            rId = sldIdLst[num_template_slides].rId
            prs.part.drop_rel(rId)
            del sldIdLst[num_template_slides]

    def _create_student_presentation(self, analytics: StudentAnalytics, student_rec: dict) -> Presentation:
        """
        Populates Slide 1 (Profile & Table) and Slide 2 (Parent Letter) for a student,
        trimming any extra slides so each student report contains exactly 2 slides.
        """
        prs = self._load_template_prs()

        # Trim extra slides if template contains previously generated slides
        if len(prs.slides) > 2:
            self._trim_to_template_slides(prs, 2)

        if len(prs.slides) == 1:
            blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
            prs.slides.add_slide(blank_layout)
        elif len(prs.slides) == 0:
            blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
            prs.slides.add_slide(blank_layout)
            prs.slides.add_slide(blank_layout)

        slide1 = prs.slides[0]
        slide2 = prs.slides[1]

        photo_path = self.image_handler.get_student_image(analytics.reg_no, analytics.name, analytics.image_filename)

        # Process placeholders & fallback logic on Slide 1
        s1_result = TemplateMapper.process_slide_placeholders(
            slide1, analytics, student_rec,
            recent_events=self.recent_events, mentor_name=self.mentor_name, mentor_phone=self.mentor_phone,
            photo_path=photo_path
        )
        self._populate_slide_1_fallback(
            slide1, analytics, student_rec, photo_path,
            photo_already_placed=s1_result.get("photo_placed", False)
        )

        # Process placeholders & fallback logic on Slide 2
        TemplateMapper.process_slide_placeholders(
            slide2, analytics, student_rec,
            recent_events=self.recent_events, mentor_name=self.mentor_name, mentor_phone=self.mentor_phone,
            photo_path=photo_path
        )
        self._populate_slide_2_fallback(slide2, analytics)

        return prs

    def _populate_slide_1_fallback(
        self,
        slide,
        analytics: StudentAnalytics,
        student_rec: dict,
        photo_path: str,
        photo_already_placed: bool = False
    ):
        """
        Fallback population for Slide 1 if explicit {{PLACEHOLDER}} tags were absent.
        Only touches the photo if it wasn't already placed via an explicit {{PHOTO}} tag,
        to avoid inserting a duplicate, overlapping photo on the slide.
        """
        for shape in slide.shapes:
            if shape.has_table:
                try:
                    TemplateMapper.populate_slide_1_table(shape.table, analytics, student_rec)
                except Exception as e:
                    tb_str = traceback.format_exc()
                    raise PPTGenerationError(
                        file_name="template_mapper.py",
                        function_name="populate_slide_1_table()",
                        student_name=analytics.name,
                        slide_name="Slide 1 (Student Profile)",
                        table_name="Academic Details Table",
                        reason=str(e),
                        traceback_str=tb_str
                    ) from e
            elif shape.has_text_frame:
                txt = shape.text_frame.text.upper()
                if any(k in txt for k in ["NAME", "1925", "REG", "ARUNPRASATH", "HARIPRIYA"]):
                    shape.text_frame.text = f"{analytics.name.upper()}, {analytics.reg_no}"

        if photo_already_placed or not photo_path or not os.path.exists(photo_path):
            return

        try:
            # Prefer the mentor's own photo placeholder shape (by size/shape),
            # so the photo sits exactly where the mentor's template design put it.
            target_shape = TemplateMapper.find_generic_photo_placeholder(slide)
            if target_shape is not None:
                left, top, width, height = target_shape.left, target_shape.top, target_shape.width, target_shape.height
                if target_shape.has_text_frame:
                    target_shape.text_frame.text = ""
                slide.shapes.add_picture(photo_path, left, top, width=width, height=height)
            else:
                # Last-resort fixed position when no placeholder shape can be detected at all.
                slide.shapes.add_picture(photo_path, Inches(0.8), Inches(1.8), width=Inches(3.2), height=Inches(3.8))
        except Exception:
            pass

    def _populate_slide_2_fallback(self, slide, analytics: StudentAnalytics):
        """
        Fallback population for Slide 2 if explicit {{PLACEHOLDER}} tags were absent.
        """
        letter_text = ParagraphBuilder.generate_parent_letter(
            analytics,
            recent_events=self.recent_events,
            mentor_name=self.mentor_name,
            mentor_phone=self.mentor_phone
        )

        letter_shape = None
        for shape in slide.shapes:
            if shape.has_text_frame:
                txt = shape.text_frame.text
                if "Dear Parents" in txt or "SIMMAM" in txt or "NPTEL" in txt:
                    letter_shape = shape
                    break

        if letter_shape and letter_shape.has_text_frame:
            letter_shape.text_frame.word_wrap = True
            letter_shape.text_frame.text = letter_text

    def _build_combined_presentation(self, individual_files: List[str], output_path: str):
        """
        Merges individual PPTX presentations into one combined cohort presentation.
        """
        if not individual_files:
            return

        shutil.copyfile(individual_files[0], output_path)
        
        if len(individual_files) > 1:
            combined = Presentation(output_path)
            for fpath in individual_files[1:]:
                try:
                    student_prs = Presentation(fpath)
                    for slide in student_prs.slides:
                        blank_layout = combined.slide_layouts[6] if len(combined.slide_layouts) > 6 else combined.slide_layouts[0]
                        new_slide = combined.slides.add_slide(blank_layout)
                        for shape in slide.shapes:
                            try:
                                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                                    # Raw XML copy only carries a relationship ID that points
                                    # into the SOURCE file's own package, not this combined one -
                                    # so the picture would render broken/missing here. Re-add it
                                    # properly instead, which copies the actual image bytes into
                                    # the combined presentation and creates a valid relationship.
                                    image_blob = shape.image.blob
                                    new_slide.shapes.add_picture(
                                        io.BytesIO(image_blob),
                                        shape.left, shape.top,
                                        width=shape.width, height=shape.height
                                    )
                                else:
                                    el = shape.element
                                    new_el = copy.deepcopy(el)
                                    new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')
                            except Exception:
                                pass
                except Exception as e:
                    print(f"Warning merging presentation '{fpath}': {e}", file=sys.stderr)

            combined.save(output_path)